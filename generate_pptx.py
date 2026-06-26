from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_slide(prs, title_text, bullets, image_path=None, notes_text=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title = slide.shapes.title
    title.text = title_text
    
    # Set bullets
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        
    # Shrink text area if image exists
    if image_path:
        body_shape.width = Inches(4.5)
        try:
            slide.shapes.add_picture(image_path, Inches(5), Inches(2), width=Inches(4.5))
        except Exception as e:
            print(f"Failed to add image {image_path}: {e}")
            
    # Add notes
    if notes_text:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

prs = Presentation()

# Slide 1: Kapak
slide_layout = prs.slide_layouts[0] # Title Slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Kubernetes Admission Webhook"
slide.placeholders[1].text = "Creation-Time Güvenlik ve Kaynak Yönetimi\nCemre Hasırcı & Taha Kulaç"
if slide.notes_slide:
    slide.notes_slide.notes_text_frame.text = "Saygıdeğer hocalarım, hoş geldiniz. Bugün sizlere Kubernetes ortamlarındaki zafiyetleri ve hatalı yapılandırmaları daha doğmadan engelleyen, bunu yaparken de bize muazzam bir yönetim arayüzü sunan Validating Admission Webhook projemizi tanıtacağız."
try:
    slide.shapes.add_picture(r"C:\Users\Taha\.gemini\antigravity\brain\34a4018b-520f-42ab-8488-4f92d0dd4f69\k8s_security_shield_1782418800363.png", Inches(7), Inches(0.5), width=Inches(2.5))
except:
    pass

# Slide 2: Projenin Kapsamı
add_slide(prs, 
          "Neleri Kontrol Ediyoruz? (Project Scope)",
          [
              "Container Security: latest tag kullanımı, Root kullanıcı, Privileged mod ve Privilege Escalation engeli.",
              "Zorunlu Kurallar: runAsNonRoot: true kullanım zorunluluğu.",
              "Storage Security: İzin verilmeyen StorageClass'ların ve Node dosya sistemine sızan hostPath kullanımının denetimi.",
              "Resource Enforcement: Tüm pod'lar için CPU ve Memory request/limit tanımlama zorunluluğu."
          ],
          r"C:\Users\Taha\.gemini\antigravity\brain\34a4018b-520f-42ab-8488-4f92d0dd4f69\k8s_vulnerability_pod_1782419397341.png",
          "Amacımız sadece tek bir güvenlik açığını kapatmak değil, Kubernetes üzerinde çalışan tüm pod'lar için aşılmaz bir standart oluşturmaktı. Bu kapsamda konteynerlerin root haklarıyla çalışmasından tutun, CPU ve RAM sınırlarının belirlenmemesine kadar cluster'ı tehlikeye atabilecek tüm senaryoları kapsayan geniş bir güvenlik ağı ördük."
)

# Slide 3: Mimari Akış
add_slide(prs,
          "Sistem Nasıl Çalışıyor? (Architecture)",
          [
              "1. Kullanıcı veya Web UI'dan Pod isteği K8s API Server'a gelir.",
              "2. API Server, bu isteği incelemesi için Admission Webhook'a iletir.",
              "3. Kurallara uygunsa ALLOW, ihlal varsa DENY kararı üretilir.",
              "4. Verilen karar ve detaylar PostgreSQL/UI üzerinden loglanır."
          ],
          r"C:\Users\Taha\.gemini\antigravity\brain\34a4018b-520f-42ab-8488-4f92d0dd4f69\k8s_webhook_shield_diagram_1782420528458.png",
          "Mimarimiz tamamen Kubernetes-native çalışmaktadır. Kullanıcı ister terminalden ister yazdığımız arayüzden bir pod göndermek istesin, bu istek API Server üzerinden bizim kural motorumuza düşer. Sistem, gelen manifestoyu saniyeler içinde analiz eder, 'Allow' veya 'Deny' kararı verir."
)

# Slide 4: Policy Matrix
add_slide(prs,
          "Ortam Bazlı Dinamik Karar Mekanizması",
          [
              "Namespace üzerindeki environment=dev veya test etiketine göre kurallar esner.",
              "Dev Ortamı: Geliştiricileri engellememek için bazı ihlaller ALLOW_WITH_WARNING olarak geçer. (Örn: hostPath kullanımı dev'de ALLOW olarak geçer.)",
              "Test/Prod Ortamı: Kurallar tavizsizdir, sıfır tolerans uygulanır. (Örn: hostPath anında DENY edilir)."
          ],
          None,
          "Projemizin en yenilikçi yanlarından biri 'Policy Matrix' yapısıdır. Node'un kendi dosyalarına erişim sağlayan tehlikeli 'hostPath' kuralı, geliştiricinin rahat çalışması için 'Dev' ortamında serbest bırakılırken, iş 'Test' veya üretim ortamına geldiğinde acımasızca reddedilir."
)

# Slide 5: Backend API ve Audit Logging
add_slide(prs,
          "FastAPI ve Güvenli Kayıt Sistemi",
          [
              "FastAPI Engine: Yüksek performanslı kural motoru. (Endpointler: /validate, /health, /audit/summary)",
              "TLS Encryption: Kubernetes API Server ile tamamen şifreli HTTPS iletişimi.",
              "PostgreSQL Entegrasyonu: Verilen tüm kararlar kalıcı Audit Log olarak yazılır.",
              "Saklanan Veriler: Namespace, Pod Adı, Karar (Allow/Deny), İhlal Edilen Kural."
          ],
          r"C:\Users\Taha\.gemini\antigravity\brain\34a4018b-520f-42ab-8488-4f92d0dd4f69\fastapi_tls_backend_1782421339886.png",
          "Sistemin kalbinde, Kubernetes API server'dan gelen istekleri karşılayan asenkron FastAPI servisimiz yatıyor. Ancak biz sadece anlık karar vermekle yetinmedik; PostgreSQL entegrasyonu sayesinde verilen her kararı, red gerekçesini saniyesi saniyesine veritabanına kaydeden tam teşekküllü bir Audit Logging sistemi kurduk."
)

# Slide 6: Web UI
add_slide(prs,
          "Modern Yönetim: Next.js Webhook Dashboard",
          [
              "Next.js & React: Webhook loglarını terminalden çıkartan modern arayüz.",
              "Single Pod Deployment: Frontend'in backend içerisinde tek bir porttan (8443) yayınlanması.",
              "Pod Config Form: Terminale dokunmadan arayüzden manifest oluşturma ve cluster'a gönderme.",
              "Log Viewer: Webhook loglarının renkli (Formatted) ve JSON olarak izlenmesi."
          ],
          r"C:\Users\Taha\Desktop\proje\k8s-admission-webhook\assets\webhook-ui-yaml.png",
          "Eğer projemizde en gurur duyduğumuz noktayı sorarsanız, o kesinlikle yazdığımız WebUI'dır. Sadece kuralları arkada çalıştıran kör bir sistem yapmak yerine, Next.js ile efsanevi bir Dashboard geliştirdik. Bu arayüz sayesinde Kubernetes terminaline hiç bulaşmadan saniyeler içinde dinamik manifestler oluşturup cluster'a gönderebiliyoruz."
)

# Slide 7: Monitoring ve Tech Stack
add_slide(prs,
          "Metrikler, İzleme ve Teknolojiler",
          [
              "Prometheus & Grafana: Webhook kararlarının ve reddedilen kural tiplerinin görselleştirilmesi.",
              "Teknoloji Yığını:",
              "- Backend: Python, FastAPI",
              "- Frontend: Next.js, React, TypeScript",
              "- Altyapı: Docker, Kubernetes, PostgreSQL"
          ],
          r"C:\Users\Taha\Desktop\proje\k8s-admission-webhook\assets\grafana-dashboard.png",
          "Son olarak, yöneticilere 'cluster'da neler oluyor?' sorusunun cevabını verebilmek için Prometheus ve Grafana entegrasyonumuzu yaptık. Python, React ve Kubernetes'in gücünü birleştirdiğimiz bu proje; salt bir güvenlik engelleyicisi değil, tam bir bulut güvenlik ürününe dönüşmüştür."
)

prs.save("Final_Sunumu.pptx")
print("Presentation generated successfully.")
